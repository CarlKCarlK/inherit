#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
import subprocess
from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pygments import lex
from pygments.lexers import RustLexer
from pygments.token import Token


def extract_markdown_parts(puzzle_md: str) -> tuple[str, str, str]:
    title_match = re.search(r"^#\s+(.+)$", puzzle_md, re.M)
    title = title_match.group(1).strip() if title_match else "Puzzle"

    mermaid_match = re.search(r"```mermaid\n([\s\S]*?)\n```", puzzle_md)
    if not mermaid_match:
        raise RuntimeError("No mermaid block found in puzzle markdown")
    mermaid = mermaid_match.group(1).strip() + "\n"

    prose = re.sub(r"^#\s+.+$", "", puzzle_md, flags=re.M)
    prose = re.sub(r"```mermaid[\s\S]*?```", "", prose).strip()
    prose = re.sub(r"\s+", " ", prose)
    return title, prose, mermaid


def resolve_mmdc_command(project_root: Path, explicit_mmdc_path: str | None) -> list[str]:
    if explicit_mmdc_path:
        return [explicit_mmdc_path]

    local_bin = project_root / "node_modules" / ".bin" / "mmdc"
    if local_bin.exists():
        return [str(local_bin)]

    system_mmdc = shutil.which("mmdc")
    if system_mmdc:
        return [system_mmdc]

    return ["npx", "@mermaid-js/mermaid-cli"]


def render_mermaid_png(
    mermaid_text: str,
    out_png: Path,
    workdir: Path,
    explicit_mmdc_path: str | None,
) -> None:
    init_block = """%%{init: {
  "theme": "base",
  "themeVariables": {
    "background": "#0B0F14",
    "primaryColor": "#111827",
    "primaryTextColor": "#F9FAFB",
    "primaryBorderColor": "#22D3EE",
    "lineColor": "#22D3EE",
    "tertiaryColor": "#0F172A"
  }
}}%%"""
    if "%%{init:" not in mermaid_text:
        mermaid_text = f"{init_block}\n{mermaid_text}"

    temp_mmd = workdir / ".tmp_mermaid.mmd"
    temp_pptr = workdir / ".tmp_puppeteer_config.json"
    temp_mmd.write_text(mermaid_text, encoding="utf-8")
    temp_pptr.write_text(
        '{"args":["--no-sandbox","--disable-setuid-sandbox"]}',
        encoding="utf-8",
    )
    try:
        cmd = resolve_mmdc_command(workdir, explicit_mmdc_path) + [
            "-i",
            str(temp_mmd),
            "-o",
            str(out_png),
            "-b",
            "transparent",
            "-p",
            str(temp_pptr),
        ]
        subprocess.run(cmd, cwd=workdir, check=True)
    finally:
        if temp_mmd.exists():
            temp_mmd.unlink()
        if temp_pptr.exists():
            temp_pptr.unlink()


def find_existing_mermaid_png(example_number: int, docs_dir: Path) -> Path | None:
    candidates = [
        docs_dir / f"{example_number}_puzzle_diagram.png",
        docs_dir / f"{example_number}_puzzle-1.png",
        docs_dir / f"{example_number}_puzzle.png",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def token_color(token_type) -> tuple[int, int, int]:
    # VS Code dark-ish palette
    if token_type in Token.Keyword or token_type in Token.Operator.Word:
        return (197, 134, 192)
    if token_type in Token.Name.Builtin or token_type in Token.Name.Class or token_type in Token.Name.Namespace or token_type in Token.Name.Type:
        return (78, 201, 176)
    if token_type in Token.Literal.String:
        return (206, 145, 120)
    if token_type in Token.Literal.Number:
        return (181, 206, 168)
    if token_type in Token.Comment:
        return (106, 153, 85)
    if token_type in Token.Name.Function:
        return (220, 220, 170)
    return (220, 220, 220)


def token_color_rgb(token_type) -> RGBColor:
    r, g, b = token_color(token_type)
    return RGBColor(r, g, b)


def load_mono_font(size: int) -> ImageFont.FreeTypeFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
    ]
    for path in candidates:
        p = Path(path)
        if p.exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def render_rust_code_png(code: str, out_png: Path, width: int = 1800) -> None:
    font = load_mono_font(24)
    pad_x = 28
    pad_y = 22
    line_gap = 6  # tighter spacing

    lines = code.splitlines()
    draw_dummy = ImageDraw.Draw(Image.new("RGB", (10, 10)))
    line_height = font.getbbox("Ag")[3] - font.getbbox("Ag")[1]
    content_height = len(lines) * line_height + max(0, len(lines) - 1) * line_gap
    height = pad_y * 2 + content_height

    image = Image.new("RGB", (width, height), (30, 30, 30))
    draw = ImageDraw.Draw(image)

    y = pad_y
    lexer = RustLexer()
    for line in lines:
        x = pad_x
        for token_type, value in lex(line, lexer):
            if not value:
                continue
            value = value.replace("\n", "")
            if not value:
                continue
            draw.text((x, y), value, font=font, fill=token_color(token_type))
            x += int(draw.textlength(value, font=font))
        y += line_height + line_gap

    image.save(out_png)


def build_ppt(
    title: str,
    prose: str,
    mermaid_png: Path,
    code_png: Path,
    code_text: str,
    example_number: int,
    out_pptx: Path,
    code_mode: str,
) -> None:
    prs = Presentation()
    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches
    bg_color = RGBColor(247, 247, 243)  # #F7F7F3
    text_color = RGBColor(31, 36, 48)   # #1F2430
    accent_color = RGBColor(91, 124, 250)  # #5B7CFA

    # Slide 1: Puzzle (edge-to-edge dark with right white rail)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    rail_w = 1.35
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(slide_w), Inches(slide_h))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(0, 0, 0)
    bg1.line.fill.background()

    rail1 = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(slide_w - rail_w),
        Inches(0),
        Inches(rail_w),
        Inches(slide_h),
    )
    rail1.fill.solid()
    rail1.fill.fore_color.rgb = RGBColor(250, 250, 250)
    rail1.line.color.rgb = RGBColor(220, 220, 220)
    rail1.line.width = Pt(0.75)

    label1 = s1.shapes.add_textbox(
        Inches(slide_w - rail_w + 0.1),
        Inches(0.3),
        Inches(rail_w - 0.2),
        Inches(slide_h - 0.6),
    )
    label1.rotation = 270
    l1tf = label1.text_frame
    l1tf.clear()
    l1tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    l1p = l1tf.paragraphs[0]
    l1p.text = f"Puzzle {example_number}"
    l1p.font.name = "Aptos"
    l1p.font.size = Pt(28)
    l1p.font.color.rgb = RGBColor(20, 20, 20)

    content_w = slide_w - rail_w
    prose_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(content_w - 1.2), Inches(1.15))
    ptf = prose_box.text_frame
    ptf.clear()
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    for part in re.split(r"(`[^`]+`)", prose):
        if not part:
            continue
        run = pp.add_run()
        if part.startswith("`") and part.endswith("`"):
            run.text = part[1:-1]
            run.font.name = "Consolas"
            run.font.color.rgb = RGBColor(255, 255, 255)
        else:
            run.text = part
            run.font.name = "Aptos"
            run.font.color.rgb = RGBColor(230, 230, 230)
        run.font.size = Pt(16)

    with Image.open(mermaid_png) as img:
        mw, mh = img.size
    ratio = mw / mh
    diagram_top = 1.85
    avail_h = slide_h - diagram_top - 0.35
    avail_w = content_w - 0.6
    pic_w = min(avail_w, avail_h * ratio)
    pic_h = pic_w / ratio
    left = (content_w - pic_w) / 2
    panel = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(max(0.1, left - 0.08)),
        Inches(max(0.1, diagram_top - 0.08)),
        Inches(pic_w + 0.16),
        Inches(pic_h + 0.16),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(11, 15, 20)  # match dark mermaid background
    panel.line.fill.background()
    s1.shapes.add_picture(str(mermaid_png), Inches(left), Inches(diagram_top), width=Inches(pic_w), height=Inches(pic_h))

    # Solution slides: all-black, edge-to-edge, no title/rail, 21 lines per slide
    code_lines = code_text.splitlines()
    lines_per_slide = 26
    for chunk_start in range(0, len(code_lines), lines_per_slide):
        chunk = code_lines[chunk_start : chunk_start + lines_per_slide]
        solution_slide = prs.slides.add_slide(prs.slide_layouts[6])
        black_bg = solution_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(slide_w), Inches(slide_h)
        )
        black_bg.fill.solid()
        black_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
        black_bg.line.fill.background()

        if code_mode == "image":
            # image mode not paginated by token image generation in this path; keep fallback
            with Image.open(code_png) as img:
                image_width, image_height = img.size
            ratio = image_width / image_height
            avail_w = slide_w
            avail_h = slide_h
            render_w = min(avail_w, avail_h * ratio)
            render_h = render_w / ratio
            solution_slide.shapes.add_picture(
                str(code_png),
                Inches((slide_w - render_w) / 2),
                Inches((slide_h - render_h) / 2),
                width=Inches(render_w),
                height=Inches(render_h),
            )
            continue

        code_box = solution_slide.shapes.add_textbox(
            Inches(0.18),
            Inches(0.12),
            Inches(slide_w - 0.36),
            Inches(slide_h - 0.24),
        )
        text_frame = code_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = False

        for line_index, line in enumerate(chunk):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.space_after = Pt(0)
            paragraph.space_before = Pt(0)
            paragraph.line_spacing = 0.95

            for token_type, value in lex(line, RustLexer()):
                value = value.replace("\n", "")
                if not value:
                    continue
                run = paragraph.add_run()
                run.text = value
                run.font.name = "Consolas"
                run.font.size = Pt(16)
                run.font.color.rgb = token_color_rgb(token_type)

    prs.save(out_pptx)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build puzzle+solution slides for one example")
    parser.add_argument("--example", type=int, required=True, help="Example number, e.g. 1")
    parser.add_argument("--project-root", type=Path, default=Path.cwd())
    parser.add_argument(
        "--mmdc-path",
        default=None,
        help="Path to Mermaid CLI binary (mmdc). If omitted, prefers ./node_modules/.bin/mmdc, then system mmdc, then npx fallback.",
    )
    parser.add_argument(
        "--code-mode",
        choices=["image", "editable"],
        default="editable",
        help="Use image-based or editable-text code rendering in the solution slide",
    )
    parser.add_argument(
        "--offline",
        action="store_true",
        help="Skip Mermaid rendering and reuse an existing diagram image from docs/",
    )
    args = parser.parse_args()

    root = args.project_root.resolve()
    examples = root / "examples"
    docs = root / "docs"
    docs.mkdir(exist_ok=True)

    puzzle_md_path = examples / f"{args.example}_puzzle.md"
    solution_rs_candidates = sorted(examples.glob(f"{args.example}_*.rs"))
    if not solution_rs_candidates:
        raise RuntimeError(f"No solution .rs file found for example {args.example}")
    solution_rs_path = solution_rs_candidates[0]

    puzzle_md = puzzle_md_path.read_text(encoding="utf-8")
    solution_code = solution_rs_path.read_text(encoding="utf-8")

    title, prose, mermaid = extract_markdown_parts(puzzle_md)

    mermaid_png = docs / f"{args.example}_puzzle_diagram.png"
    code_png = docs / f"{args.example}_solution_code.png"
    out_pptx = docs / f"example{args.example}_walkthrough.pptx"

    if args.offline:
        existing = find_existing_mermaid_png(args.example, docs)
        if not existing:
            raise RuntimeError(
                "Offline mode requested, but no existing diagram image found in docs/"
            )
        if existing != mermaid_png:
            mermaid_png.write_bytes(existing.read_bytes())
    else:
        render_mermaid_png(mermaid, mermaid_png, root, args.mmdc_path)
    render_rust_code_png(solution_code, code_png)
    build_ppt(
        title,
        prose,
        mermaid_png,
        code_png,
        solution_code,
        args.example,
        out_pptx,
        args.code_mode,
    )

    print(out_pptx)


if __name__ == "__main__":
    main()
